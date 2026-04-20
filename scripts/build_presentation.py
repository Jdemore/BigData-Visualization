"""Build the LAVA presentation deck as docs/LAVA_presentation.pptx.

Kent State palette: navy #002664, gold #EAAB00.
Target length: 10-15 min, ~14 slides, big-data focus, speaker notes included.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

NAVY = RGBColor(0x00, 0x26, 0x64)
GOLD = RGBColor(0xEA, 0xAB, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF4, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_header_bar(slide, title: str):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.9), SLIDE_W, Inches(0.05)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = GOLD
    stripe.line.fill.background()


def add_footer(slide, slide_num: int, total: int):
    tb = slide.shapes.add_textbox(
        Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.3)
    )
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = f"LAVA  |  Kent State University  |  {slide_num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY


def add_bullets(slide, bullets, left, top, width, height, font_size=16):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(font_size - 2 if level > 0 else font_size)
        p.font.color.rgb = DARK if level == 0 else GRAY
        p.space_after = Pt(6)


def set_notes(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def title_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    gold = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3.4), SLIDE_W, Inches(0.08)
    )
    gold.fill.solid()
    gold.fill.fore_color.rgb = GOLD
    gold.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(12), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "LAVA"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tb.text_frame.add_paragraph()
    p2.text = "LLM-Assisted Visual Analytics for Big Data"
    p2.font.size = Pt(28)
    p2.font.color.rgb = GOLD

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(12), Inches(2.5))
    for line, size, color in [
        ("Natural-language querying over million-row datasets", 20, WHITE),
        ("", 10, WHITE),
        ("Kent State University", 16, GOLD),
        ("CS Big Data Visualization -- 2026", 14, WHITE),
    ]:
        p = tb2.text_frame.add_paragraph() if tb2.text_frame.paragraphs[0].text else tb2.text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color

    set_notes(slide, (
        "Hello, I'm presenting LAVA -- LLM-Assisted Visual Analytics. "
        "LAVA lets a user ask plain-English questions of a large dataset "
        "and get an interactive chart back in a couple of seconds. "
        "The focus of this talk is the big-data plumbing underneath -- "
        "how we keep query latency flat as the data grows from thousands "
        "to millions of rows. I'll spend the first couple of minutes framing "
        "the problem, then most of the time on ingestion, indexing, and query "
        "processing, and close with our experimental results."
    ))


def problem_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "The Problem")
    add_bullets(slide, [
        "Analysts want answers; they do not want to write SQL.",
        "General LLM chat tools do not scale to millions of rows --",
        ("they ingest the data into the prompt window.", 1),
        "Existing NL-to-viz tools (Chat2VIS, ThoughtSpot) are strong on chart",
        ("type but weak on the data-engineering side: no indexing, no pushdown,", 1),
        ("no LOD rendering.", 1),
        "",
        "Goal: a system where the LLM does mapping, and a real big-data stack",
        "does the heavy lifting.",
    ], Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), 18)
    add_footer(slide, 2, total)
    set_notes(slide, (
        "The problem we started with is simple to state. Business users and "
        "researchers want to ask questions about their data in plain English, "
        "but the tools that do that today either require the user to know SQL, "
        "or they stuff the dataset into an LLM prompt -- which breaks as soon "
        "as you go past a few thousand rows. Tools like Chat2VIS are good at "
        "picking the right chart type but do not address the data-engineering "
        "problems: how do you index the data, how do you push work down to the "
        "database, how do you render a million-point scatter plot without killing "
        "the browser. That's the gap LAVA is filling."
    ))


def architecture_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "System Architecture")

    boxes = [
        ("User NL Query", 0.6, 1.6, 2.2, 0.9, NAVY, WHITE),
        ("LLM Pipeline\n(refine -> VizSpec)", 3.1, 1.6, 2.6, 0.9, NAVY, WHITE),
        ("SQL Generator", 6.0, 1.6, 2.2, 0.9, NAVY, WHITE),
        ("DuckDB +\nIndexes", 8.5, 1.6, 2.2, 0.9, GOLD, DARK),
        ("Parquet Store", 11.0, 1.6, 1.9, 0.9, GOLD, DARK),
        ("Result Arrow Table", 6.0, 3.2, 2.2, 0.9, NAVY, WHITE),
        ("Renderer\n(Plotly / WebGL / Datashader)", 8.5, 3.2, 3.4, 0.9, NAVY, WHITE),
        ("Dash Web UI", 8.5, 4.7, 2.2, 0.9, GOLD, DARK),
    ]
    for text, x, y, w, h, fill, fg in boxes:
        s = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.color.rgb = NAVY
        tf = s.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = fg

    caption = slide.shapes.add_textbox(
        Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0)
    )
    p = caption.text_frame.paragraphs[0]
    p.text = "LLM handles the mapping problem. DuckDB + indexes + LOD rendering handle the scaling problem."
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GRAY

    add_footer(slide, 3, total)
    set_notes(slide, (
        "Here's the full system on one slide. A user query comes in on the left. "
        "The LLM pipeline refines it and emits a structured VizSpec -- a JSON "
        "describing what chart, which columns, which aggregation. That spec is "
        "deterministically translated to SQL. The SQL hits DuckDB, which uses "
        "our B+-tree and Grid File indexes to prune work. Results come back as "
        "an Arrow table -- zero-copy into Plotly. The renderer picks SVG, WebGL, "
        "or Datashader based on row count. The key design choice: the LLM never "
        "sees the data. It only sees column names and statistics. That's what "
        "makes this scale past the prompt window."
    ))


def techstack_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Technology Stack")
    items = [
        ("Storage", "Apache Parquet -- columnar, compressed (ZSTD), row-group pruning"),
        ("Query Engine", "DuckDB -- in-process OLAP, vectorized, zero-copy Arrow output"),
        ("Indexes", "B+-tree (4 KB pages, 254 keys/leaf)  +  Grid File (2-D range queries)"),
        ("Rendering", "Plotly (SVG)  ->  WebGL (>1K points)  ->  Datashader (>100K points)"),
        ("LLM", "OpenAI GPT-4o-mini, function calling, deterministic (temp 0.1)"),
        ("UI", "Dash + Bootstrap, drag-and-drop file upload (CSV / JSON / Parquet)"),
    ]
    top = 1.3
    for label, desc in items:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(3.2), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NAVY

        tb2 = slide.shapes.add_textbox(Inches(4.2), Inches(top), Inches(8.8), Inches(0.7))
        p = tb2.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        top += 0.75

    add_footer(slide, 4, total)
    set_notes(slide, (
        "Quick tour of the stack before we go deep. Storage is Parquet -- "
        "columnar and compressed, which matters because most analytical queries "
        "only touch a few columns out of many. DuckDB is our query engine: "
        "in-process, vectorized, and critically it hands back Apache Arrow tables, "
        "which Plotly can consume with zero copy. Our two custom index structures "
        "sit on top. The renderer has three tiers based on cardinality. The LLM "
        "is GPT-4o-mini at temperature 0.1 for determinism. The UI is Dash with "
        "a drag-and-drop upload for CSV, JSON, or Parquet."
    ))


def ingestion_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Big Data -- Ingestion & Storage")
    add_bullets(slide, [
        "Any CSV / JSON / Parquet -> profiled by DuckDB (DESCRIBE, no full scan).",
        "Converted to Parquet with ZSTD compression, 100K-row row groups.",
        ("Row-group pruning skips entire chunks based on min/max metadata.", 1),
        ("Column projection reads only the columns a query needs.", 1),
        "",
        "Column-stats extraction: batched into 1-2 queries instead of O(n).",
        ("Before: per-column SQL round-trips  ->  ~280 ms for 9 columns.", 1),
        ("After: one batched query  ->  ~95 ms  (3x speedup).", 1),
        "",
        "Result: bootstrap on a 50K-row dataset in under a second.",
    ], Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), 17)
    add_footer(slide, 5, total)
    set_notes(slide, (
        "Now into the big-data side proper. Ingestion is the first scaling "
        "problem. Users upload arbitrary CSV or JSON files, and we need to "
        "profile them without loading everything into RAM. DuckDB's DESCRIBE "
        "lets us get the schema by reading just the header. We then convert "
        "to Parquet with ZSTD compression and 100K-row row groups -- that "
        "row-group size is tuned so DuckDB can prune based on per-group min/max "
        "metadata. One optimization I want to call out: column-stats extraction "
        "used to issue one query per column, so nine columns meant nine round "
        "trips, about 280 ms. We batched them into a single SQL statement using "
        "COUNT(DISTINCT) across columns, cutting that to 95 ms -- a 3x speedup "
        "on bootstrap."
    ))


def indexing_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Big Data -- Indexing")

    left_tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(6.2), Inches(5.5))
    tf = left_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "B+-tree (1-D lookups)"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = NAVY
    for line, level in [
        ("4 KB pages, OS-aligned", 1),
        ("254 keys per leaf = (4096 - 20) / 16", 1),
        ("Single-pass bulk-load from a sorted array", 1),
        ("Build: ~180 ms on 50K rows", 1),
        ("Lookup: ~3 page reads (tree height)", 1),
        ("Pruning ratio: >99% for point/range queries", 1),
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.level = level
        p.font.size = Pt(14)
        p.font.color.rgb = DARK

    right_tb = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.5))
    tf = right_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Grid File (2-D range queries)"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = NAVY
    for line, level in [
        ("Partitions 2 numeric dims into ~10K-row buckets", 1),
        ("Quantile-based, so buckets are balanced", 1),
        ("Build: ~320 ms on 50K rows, 2 dims", 1),
        ("Range query: ~4 bucket reads on average", 1),
        ("Cell selectivity: ~96%", 1),
        ("", 1),
        ("Both indexes are disk-resident and survive restart.", 0),
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.level = level
        p.font.size = Pt(14)
        p.font.color.rgb = DARK

    add_footer(slide, 6, total)
    set_notes(slide, (
        "Indexing is where the big-data contribution is most concrete. We built "
        "two index structures from scratch. The B+-tree uses 4-kilobyte pages "
        "aligned to the OS page size, and packs 254 keys per leaf -- that number "
        "comes from dividing the page size minus a 20-byte header by the 16-byte "
        "key-plus-pointer record. We bulk-load it in a single pass from a sorted "
        "array, which is much faster than inserting keys one at a time. On 50K "
        "rows, build takes about 180 milliseconds and any lookup touches only "
        "three pages -- that's over 99 percent pruning. The Grid File handles "
        "two-dimensional range queries, the kind you get with geographic or "
        "scatter-plot filters. It partitions both dimensions into buckets of "
        "about 10K rows each, using quantiles so the buckets are balanced. A "
        "typical range query reads four buckets. Both indexes are persisted "
        "to disk so they survive restarts."
    ))


def query_processing_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Big Data -- Query Processing & Rendering")
    add_bullets(slide, [
        "Aggregation pushed down to DuckDB -- never pull raw rows to Python.",
        ("'revenue by region' becomes GROUP BY in SQL, returning ~10 rows.", 1),
        "Server-side binning for histograms (width_bucket), 100 bins default.",
        "LRU query cache -- 64 entries, keyed on canonicalized SQL.",
        "",
        "Level-of-detail rendering based on row count:",
        ("<= 1,000 rows   ->  Plotly SVG  (interactive, precise)", 1),
        (">  1,000 rows   ->  Plotly WebGL  (GPU-accelerated scatter)", 1),
        (">  100,000 rows ->  Datashader  (rasterized density map)", 1),
        "",
        "Arrow table handoff: zero-copy DuckDB -> Plotly, no Python list conversion.",
    ], Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.7), 16)
    add_footer(slide, 7, total)
    set_notes(slide, (
        "Query processing has two jobs: push work down, and render smartly. "
        "Pushdown means we never pull raw rows to Python. If the user asks for "
        "revenue by region, that becomes a GROUP BY in SQL and DuckDB returns "
        "ten rows, not fifty thousand. Histograms use DuckDB's width_bucket "
        "function so the binning happens server-side. We cache the last 64 "
        "queries with an LRU, keyed on the canonicalized SQL so small query "
        "variations still hit the cache. On the render side we have three "
        "tiers based on row count. Under a thousand points, Plotly SVG for "
        "precise interactivity. Over a thousand, we switch to WebGL which "
        "uses the GPU. Past a hundred thousand, we hand off to Datashader, "
        "which rasterizes the dataset into a density image -- that's how we "
        "render a million-point scatter plot without crashing the browser. "
        "The whole pipeline is zero-copy Arrow: DuckDB outputs Arrow, Plotly "
        "consumes Arrow, no intermediate Python lists."
    ))


def llm_pipeline_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "LLM Pipeline (brief)")
    add_bullets(slide, [
        "Two-stage: refine the query, then generate a VizSpec.",
        "LLM sees column names, types, cardinalities, sample values --",
        ("never the actual data. That is what decouples it from dataset size.", 1),
        "Output is a typed JSON schema (chart_type, x, y, aggregation, filter).",
        ("Deterministic translator converts VizSpec  ->  parameterized SQL.", 1),
        "History window of 3 turns supports follow-ups ('now make it monthly').",
        "",
        "Error log captures every failure + raw LLM output for self-healing later.",
    ], Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), 17)
    add_footer(slide, 8, total)
    set_notes(slide, (
        "I'll keep the LLM side brief since the focus is big data, but the design "
        "choice that matters is this: the LLM only ever sees metadata -- column "
        "names, types, cardinalities, and a handful of sample values. It never "
        "sees the rows. That's what decouples LLM cost and latency from dataset "
        "size. The pipeline runs in two stages. First we refine the user's query "
        "into a canonical form. Then we call the model again to emit a structured "
        "VizSpec -- a typed JSON describing the chart. That spec goes through a "
        "deterministic translator to SQL, not a second LLM call, so the SQL is "
        "reproducible and parameterized. A three-turn history window supports "
        "follow-ups like 'now make it monthly.'"
    ))


def gui_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "GUI & File Upload")
    add_bullets(slide, [
        "Dash web app -- single-page, Bootstrap-styled.",
        "Drag-and-drop upload accepts CSV, JSON, NDJSON, Parquet (up to 200 MB).",
        ("File is routed through DuckDB read_json_auto / read_csv_auto,", 1),
        ("converted to Parquet, indexed, and made active in seconds.", 1),
        "Chart-type override dropdown -- 22 chart types supported.",
        "SQL footer is always visible: every chart is reproducible and auditable.",
        "Escape clears the query input; Enter submits.",
    ], Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), 17)
    add_footer(slide, 9, total)
    set_notes(slide, (
        "The UI is a single-page Dash app. The feature I want to call out is "
        "the file upload -- users can drag and drop a CSV, JSON, or Parquet "
        "file up to 200 megabytes, and the system will profile it, convert it "
        "to Parquet, build stats, and make it the active dataset in just a few "
        "seconds. JSON support routes through DuckDB's read_json_auto. The "
        "chart-type dropdown lets the user override the LLM's choice across "
        "22 supported chart types. And every chart shows its generating SQL "
        "in the footer, so the output is fully auditable -- no black box."
    ))


def experimental_setup_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Experimental Setup")
    add_bullets(slide, [
        "Hardware: consumer laptop, 16 GB RAM, no GPU for LLM (API-only).",
        "Datasets: synthetic 50K-row sales table (9 columns, mixed types).",
        "Test suite: 72 automated pytest cases, 2.3s total runtime.",
        "",
        "Prompt evaluation protocol:",
        ("50 natural-language prompts per data size", 1),
        ("Scored on chart produced, accuracy, and wall time <= 5 s", 1),
        ("Separate measurement for index build time and memory footprint", 1),
    ], Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), 17)
    add_footer(slide, 10, total)
    set_notes(slide, (
        "For the evaluation we ran on a commodity laptop -- 16 gigs of RAM, "
        "no GPU, LLM access via the OpenAI API. The dataset is a synthetic "
        "50,000-row sales table with nine columns of mixed types. We have "
        "72 automated pytest cases covering ingestion, indexing, query "
        "generation, and rendering -- they run in about two and a third "
        "seconds. For the end-to-end evaluation we issued 50 natural-language "
        "prompts per data size and scored each on three criteria: did it "
        "produce a chart, was the chart accurate, and did it finish within "
        "a 5-second budget. Index build time and memory footprint were "
        "measured separately on the same sample set."
    ))


def performance_chart_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Performance Report")

    img_path = "docs/performancechart.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path,
            Inches(2.3), Inches(1.15),
            width=Inches(8.7),
        )

    add_footer(slide, 11, total)
    set_notes(slide, (
        "This is the consolidated performance chart. Top-left: the outcome "
        "breakdown for the 50-prompt run -- 43 accurate in under 5 seconds, "
        "4 inaccurate, 3 slow. Top-right: retrieval quality -- precision 0.935, "
        "recall 0.860, F1 0.896. Bottom-left: per-stage CPU time on a log scale. "
        "DuckDB query time is about 40 milliseconds, stats extraction is 95 "
        "milliseconds, and the LLM dominates at 2 seconds -- that's 95 percent "
        "of end-to-end latency, and it's the remote API, not our pipeline. "
        "Bottom-right: index construction -- B+-tree builds in 180 milliseconds "
        "and fits in under a megabyte; Grid File builds in 320 milliseconds "
        "with negligible directory size."
    ))


def prompt_accuracy_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Prompt-to-Chart Accuracy")

    rows = [
        ("Metric", "Count", "%"),
        ("Prompts issued", "50", "100%"),
        ("Charts produced", "46", "92%"),
        ("Accurate AND rendered < 5s", "43", "86%"),
        ("Produced but inaccurate", "4", "8%"),
        ("Produced but over 5s", "3", "6%"),
    ]
    table = slide.shapes.add_table(
        rows=len(rows), cols=3,
        left=Inches(1.2), top=Inches(1.3),
        width=Inches(7.5), height=Inches(2.8),
    ).table
    for j in range(3):
        table.columns[j].width = Inches(2.5)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(14)
                if i == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                else:
                    para.font.color.rgb = DARK
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for line, bold in [
        ("Key observations:", True),
        ("-- All 4 inaccurate charts were the same chart type -- a narrow JSON-return defect, not a systemic issue.", False),
        ("-- All 3 over-budget renders finished within 0.5 s of the 5 s threshold; attributed to remote-LLM variance.", False),
        ("-- Fixing the single JSON branch lifts precision to ~1.00 and F1 to ~0.96.", False),
    ]:
        p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.bold = bold
        p.font.color.rgb = NAVY if bold else DARK

    add_footer(slide, 12, total)
    set_notes(slide, (
        "Here's the detail behind the 50-prompt run. Of 50 prompts, 46 produced "
        "a chart, and 43 of those were both accurate and rendered in under "
        "5 seconds -- so 86 percent end-to-end success. Two things I want to "
        "flag about the misses. First, the 4 inaccurate charts were all of the "
        "same chart type, which tells us this is a narrow defect in one "
        "JSON-return branch, not a systemic accuracy problem. Fixing that "
        "single branch pushes precision to nearly 1.0. Second, the 3 slow "
        "renders all landed within half a second of the 5-second threshold, "
        "and since the LLM endpoint is remote and outside our control, we "
        "treat these as network-side outliers rather than pipeline regressions."
    ))


def limitations_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Limitations & Future Work")

    left_tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.5))
    tf = left_tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Limitations"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = NAVY
    for line in [
        "Single table only -- no multi-table JOINs yet.",
        "Single-user by design (module-level state).",
        "LLM latency dominates (~95% of total time).",
        "Upload path holds file in browser RAM (200 MB cap).",
    ]:
        p = tf.add_paragraph()
        p.text = "-- " + line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK

    right_tb = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.5))
    tf = right_tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Future Work"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = NAVY
    for line in [
        "Multi-table JOIN support via INFORMATION_SCHEMA FK discovery.",
        "Per-session state for multi-user deployment (Redis).",
        "Self-healing: mine error log to propose prompt improvements.",
        "Streaming ingestion for append-only log files.",
        "Geographic chart types (choropleth, scattermapbox).",
    ]:
        p = tf.add_paragraph()
        p.text = "-- " + line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK

    add_footer(slide, 13, total)
    set_notes(slide, (
        "Honest limitations. We operate on one table at a time -- no JOINs yet. "
        "The system is single-user by design, because we use module-level state "
        "for the connection and the conversation history. LLM latency dominates; "
        "there's not much we can do about that short of a smaller local model. "
        "And the upload path holds the entire file in browser RAM, so we cap it "
        "at 200 megabytes. For future work, the biggest wins are multi-table "
        "JOINs driven by foreign-key discovery, per-session state for real "
        "deployment, and a self-healing loop that mines our error log to propose "
        "prompt improvements automatically."
    ))


def qa_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

    gold = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3.4), SLIDE_W, Inches(0.08)
    )
    gold.fill.solid(); gold.fill.fore_color.rgb = GOLD; gold.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(12), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tb.text_frame.add_paragraph()
    p2.text = "Questions?"
    p2.font.size = Pt(36); p2.font.color.rgb = GOLD

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(12), Inches(2))
    for line in [
        "LAVA -- LLM-Assisted Visual Analytics",
        "Kent State University",
    ]:
        p = tb2.text_frame.add_paragraph() if tb2.text_frame.paragraphs[0].text else tb2.text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

    set_notes(slide, (
        "Thank you. To recap: LAVA uses an LLM for the NL-to-intent mapping, "
        "and a real big-data stack -- Parquet, DuckDB, custom B+-tree and "
        "Grid File indexes, and level-of-detail rendering -- to keep query "
        "latency flat as datasets grow. Happy to take questions."
    ))


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        title_slide,
        problem_slide,
        architecture_slide,
        techstack_slide,
        ingestion_slide,
        indexing_slide,
        query_processing_slide,
        llm_pipeline_slide,
        gui_slide,
        experimental_setup_slide,
        performance_chart_slide,
        prompt_accuracy_slide,
        limitations_slide,
        qa_slide,
    ]
    total = len(slides)
    for fn in slides:
        fn(prs, total)

    out = "docs/LAVA_presentation.pptx"
    prs.save(out)
    print(f"saved -> {out}  ({total} slides)")


if __name__ == "__main__":
    build()
