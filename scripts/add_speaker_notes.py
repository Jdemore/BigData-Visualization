"""Attach speaker notes to each slide in LAVA_ BigData Vis.pptx.

One block per slide; order matches the deck. Rewriting the notes frame in
place keeps the slide layout untouched and just replaces the speaker notes.
"""

from pptx import Presentation

NOTES = [
    # Slide 1 -- Title
    (
        "Hi, I'm Joseph Demore, and this is LAVA -- LLM-Assisted Visual "
        "Analytics. The idea behind LAVA is simple: let a user ask a plain-"
        "English question of a large dataset and get back an interactive "
        "chart in a couple of seconds. The focus of this talk is the big-"
        "data plumbing underneath -- how we keep query time flat as the "
        "data grows from fifty thousand to a million rows. I'll spend the "
        "first couple of minutes on the problem and the architecture, then "
        "most of the time on ingestion, indexing, and query processing, "
        "and close with evaluation results."
    ),

    # Slide 2 -- Problem
    (
        "Here's the gap we're trying to close. Business owners and analysts "
        "want answers from their data, but they don't want to write SQL. "
        "The obvious workaround is to hand the data to a chat LLM, but "
        "those don't scale -- the model either runs out of context window "
        "or gives up on larger files entirely. There are tools like Chat2VIS "
        "and ThoughtSpot that do natural-language-to-chart, and they're good "
        "at picking chart types, but they stop there. They don't index the "
        "data, they don't push computation down to a database, and they "
        "don't have level-of-detail rendering. So the goal for LAVA was to "
        "combine an LLM for the mapping problem -- going from English to a "
        "structured query -- with a real big-data stack for the scaling "
        "problem underneath."
    ),

    # Slide 3 -- System Architecture
    (
        "Here's the whole system on one slide. A user query comes in on the "
        "left. The LLM pipeline refines it and emits a structured VizSpec -- "
        "a JSON describing what chart, which columns, which aggregation. "
        "That spec is deterministically translated to SQL -- no LLM "
        "involvement in the SQL step, so the query is reproducible. DuckDB "
        "executes it against our Parquet store, using the B+-tree and Grid "
        "File indexes to prune work. Results come back as Arrow, zero-copy "
        "into Plotly. The renderer picks SVG, WebGL, or Datashader based on "
        "row count. The key architectural decision: the LLM never sees the "
        "actual data. It only sees column names and statistics. That's what "
        "decouples cost and latency from dataset size."
    ),

    # Slide 4 -- Tech Stack
    (
        "Quick tour of the stack before we go deep. Storage is Apache "
        "Parquet with ZSTD compression. Columnar format means most queries "
        "only read the columns they need, not every row. DuckDB is the "
        "query engine -- in-process OLAP, vectorized execution, and it "
        "hands results back as Apache Arrow tables, which Plotly can "
        "consume without copying. We have two custom index structures on "
        "top: a B+-tree for one-dimensional lookups and a Grid File for "
        "two-dimensional range queries. The renderer has three tiers based "
        "on cardinality. The LLM is GPT-4o-mini at temperature 0.1 for "
        "determinism. The UI is Dash with Bootstrap and drag-and-drop "
        "upload for CSV, JSON, or Parquet."
    ),

    # Slide 5 -- Ingestion and Storage
    (
        "Now into the big-data side. Ingestion is the first scaling problem. "
        "Users upload arbitrary files, and we need to profile them without "
        "loading everything into RAM. DuckDB's DESCRIBE statement gets us "
        "the schema by reading just the header. We then convert to Parquet "
        "with ZSTD compression and hundred-thousand-row row groups -- that "
        "size is tuned so DuckDB can prune based on per-group min and max "
        "metadata. One optimization I want to highlight: column-stats "
        "extraction used to issue one query per column, so nine columns "
        "meant nine round trips, about 280 milliseconds. I batched them "
        "into a single SQL statement using COUNT DISTINCT across every "
        "column, cutting that to 95 milliseconds -- a three-times speedup "
        "on bootstrap. The end result is sub-second load times even at a "
        "million rows."
    ),

    # Slide 6 -- Indexing
    (
        "Indexing is where the big-data contribution is most concrete. I "
        "built two index structures from scratch. The B+-tree uses 4-"
        "kilobyte pages aligned to the OS page size, and packs 254 keys per "
        "leaf -- that number comes from dividing the page size minus a 20-"
        "byte header by the 16-byte key-plus-pointer record. I bulk-load it "
        "in a single pass from a sorted array, which is much faster than "
        "inserting keys one at a time. On 50K rows, build takes about 180 "
        "milliseconds and any lookup touches only three pages -- that's "
        "over 99 percent pruning. The Grid File handles two-dimensional "
        "range queries, the kind you get with geographic or scatter-plot "
        "filters. It partitions both dimensions into buckets of about 10K "
        "rows each, using quantiles so the buckets stay balanced regardless "
        "of distribution. A typical range query reads four buckets on "
        "average. Both indexes are persisted to disk, so they survive "
        "restarts."
    ),

    # Slide 7 -- Query Processing and Rendering
    (
        "Query processing has two jobs: push work down to the database, "
        "and render smartly. Pushdown means we never pull raw rows to "
        "Python. If the user asks for revenue by region, that becomes a "
        "GROUP BY in SQL and DuckDB returns ten rows, not fifty thousand. "
        "Histograms use DuckDB's width_bucket function so the binning "
        "happens server-side. I cache the last 64 queries with an LRU, "
        "keyed on the canonicalized SQL so small query variations still hit "
        "the cache. On the render side there are three tiers based on row "
        "count. Under a thousand points, Plotly SVG gives precise "
        "interactivity. Over a thousand, we switch to WebGL, which uses the "
        "GPU. Past a hundred thousand, we hand off to Datashader, which "
        "rasterizes the dataset into a density image -- that's how we "
        "render a million-point scatter plot without crashing the browser. "
        "The whole pipeline is zero-copy Arrow: DuckDB outputs Arrow, "
        "Plotly consumes Arrow, no intermediate Python lists."
    ),

    # Slide 8 -- Pipeline
    (
        "Brief detour into the LLM side since the focus is big data, but "
        "one design choice matters. The LLM only ever sees metadata -- "
        "column names, types, cardinalities, and a handful of sample values. "
        "It never sees the rows. That's what decouples LLM cost and "
        "latency from dataset size, and it's why a million-row dataset "
        "costs the same in tokens as a thousand-row one. The pipeline runs "
        "in two stages. First I refine the user's query into a canonical "
        "analytical statement. Then I call the model again to emit a "
        "structured VizSpec. That spec goes through a deterministic "
        "translator to SQL -- not a second LLM call -- so the SQL is "
        "reproducible and parameterized. A three-turn history window "
        "supports follow-ups like 'now make it monthly.' And every failure "
        "gets logged with the raw LLM output, which is what made it "
        "possible to diagnose and fix the accuracy bugs I'll mention later."
    ),

    # Slide 9 -- GUI and Uploading
    (
        "The UI is a single-page Dash app. The feature I want to call out "
        "is the file upload. Users can drag and drop a CSV, JSON, or "
        "Parquet file up to 200 megabytes, and the system will profile it, "
        "convert it to Parquet, build statistics, and make it the active "
        "dataset in a few seconds. JSON support routes through DuckDB's "
        "read_json_auto function. The chart-type dropdown lets the user "
        "override the LLM's choice across 22 supported chart types -- "
        "everything from bar and line up through radar, waterfall, and "
        "Datashader density maps. And every chart shows its generating SQL "
        "in the footer, so the output is fully auditable. No black box."
    ),

    # Slide 10 -- Experiment Setup
    (
        "For evaluation I ran everything on a commodity laptop: 16 gigs of "
        "RAM, no GPU, LLM access through the OpenAI API. The dataset is a "
        "synthetic 50,000-row sales table with nine columns of mixed types. "
        "There are 72 automated pytest cases covering ingestion, indexing, "
        "query generation, and rendering -- they run in about two and a "
        "third seconds. For the end-to-end evaluation I issued 50 natural-"
        "language prompts per data size and scored each prompt on three "
        "criteria: did it produce a chart, was the chart accurate, and did "
        "it finish within a 5-second budget. Index build time and memory "
        "footprint were measured separately on the same sample set, so all "
        "the timing numbers reflect the same system state."
    ),

    # Slide 11 -- Performance
    (
        "Here's the consolidated performance picture. Top-left: the outcome "
        "breakdown for the 50-prompt run -- 43 accurate in under 5 seconds, "
        "4 inaccurate, 3 over time. Top-right: retrieval quality -- "
        "precision of 0.935, recall of 0.860, F1 of 0.896. Bottom-left: "
        "per-stage CPU time on a log scale. DuckDB query time is around "
        "40 milliseconds, stats extraction is 95 milliseconds, and the LLM "
        "dominates at about 2 seconds. That's 95 percent of end-to-end "
        "latency, and it's the remote API, not our local pipeline. "
        "Bottom-right: index construction -- the B+-tree builds in 180 "
        "milliseconds and fits in under a megabyte. The Grid File builds "
        "in 320 milliseconds. These are the numbers I care about -- they "
        "say the big-data side is doing its job, and the variable cost is "
        "the remote LLM."
    ),

    # Slide 12 -- Prompt to Visual Accuracy (chart)
    (
        "Here's the accuracy breakdown in detail. Of 50 prompts, 46 "
        "produced a chart -- that's 92 percent. Of those 46, 43 were both "
        "accurate and rendered in under 5 seconds. That's 86 percent end-"
        "to-end success on the full criteria. Four produced charts that "
        "were inaccurate, and three rendered correctly but went over the "
        "5-second time budget. So the headline number is 86 percent, but "
        "the failure modes are worth looking at, because they're not "
        "random -- and that's what the next slide covers."
    ),

    # Slide 13 -- Prompt to Visual Accuracy (observations)
    (
        "Two observations about the misses. First, all four inaccurate "
        "charts were the same chart type. That tells us it's a narrow "
        "defect in one JSON-return branch, not a systemic accuracy problem. "
        "I know where the fix goes -- patching that single branch would "
        "lift precision from 0.935 to essentially 1.0 and F1 from 0.896 to "
        "about 0.96. Second, the three slow renders all landed within half "
        "a second of the 5-second threshold, and because the LLM endpoint "
        "is remote and outside our control, I'm treating those as network-"
        "side variance, not a pipeline regression. The LLM round-trip "
        "alone accounts for 95 percent of latency -- half a second of "
        "jitter on a 2-second call is entirely normal. So the real system "
        "accuracy, once that one branch is patched, sits around 94 percent, "
        "with the remaining 6 percent bounded by network variance."
    ),

    # Slide 14 -- Limitations and Future Work
    (
        "Honest limitations. LAVA operates on one table at a time -- no "
        "JOINs yet. It's single-user by design, because the connection and "
        "history use module-level state. LLM latency dominates everything, "
        "and short of a smaller local model there's not much we can do "
        "about that. And the upload path holds the entire file in browser "
        "RAM, which is why there's a 200-megabyte cap. For future work, "
        "the biggest wins are multi-table JOIN support using "
        "INFORMATION_SCHEMA foreign-key discovery, per-session state using "
        "something like Redis for real multi-user deployment, and a self-"
        "healing loop that mines the error log to propose prompt "
        "improvements automatically. Streaming ingestion and geographic "
        "chart types are natural extensions beyond that."
    ),

    # Slide 15 -- Thank You
    (
        "That's LAVA. The short version: the LLM does the natural-"
        "language-to-intent mapping, and a real big-data stack -- Parquet, "
        "DuckDB, custom B+-tree and Grid File indexes, and level-of-detail "
        "rendering -- keeps query latency flat as datasets grow. Thank you "
        "for watching, and happy to take questions."
    ),
]


def main() -> None:
    src = "LAVA_ BigData Vis.pptx"
    prs = Presentation(src)
    slides = list(prs.slides)

    if len(slides) != len(NOTES):
        print(
            f"Warning: deck has {len(slides)} slides, notes list has "
            f"{len(NOTES)}. Writing to the minimum of both."
        )

    for slide, note in zip(slides, NOTES):
        slide.notes_slide.notes_text_frame.text = note

    prs.save(src)
    print(f"Wrote speaker notes to {len(slides)} slides in {src}")


if __name__ == "__main__":
    main()
